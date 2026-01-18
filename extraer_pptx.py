"""
Script para extraer el contenido de texto de archivos PowerPoint (.pptx)
y guardarlo en archivos de texto (.txt) accesibles.
"""

from pptx import Presentation
from pathlib import Path
import os

def extraer_texto_pptx(ruta_pptx: str, ruta_salida: str = None) -> str:
    """
    Extrae todo el texto de un archivo PowerPoint.
    
    Args:
        ruta_pptx: Ruta al archivo .pptx
        ruta_salida: Ruta donde guardar el archivo .txt (opcional)
    
    Returns:
        El texto extraído de la presentación
    """
    prs = Presentation(ruta_pptx)
    texto_completo = []
    
    nombre_archivo = Path(ruta_pptx).stem
    texto_completo.append(f"{'='*60}")
    texto_completo.append(f"PRESENTACIÓN: {nombre_archivo}")
    texto_completo.append(f"{'='*60}\n")
    
    for num_diapositiva, diapositiva in enumerate(prs.slides, 1):
        texto_completo.append(f"\n{'─'*40}")
        texto_completo.append(f"DIAPOSITIVA {num_diapositiva}")
        texto_completo.append(f"{'─'*40}\n")
        
        for forma in diapositiva.shapes:
            if hasattr(forma, "text") and forma.text.strip():
                texto_completo.append(forma.text.strip())
                texto_completo.append("")  # Línea en blanco entre elementos
            
            # Extraer texto de tablas si existen
            if forma.has_table:
                tabla = forma.table
                texto_completo.append("[TABLA]")
                for fila in tabla.rows:
                    celdas = [celda.text.strip() for celda in fila.cells]
                    texto_completo.append(" | ".join(celdas))
                texto_completo.append("")
    
    resultado = "\n".join(texto_completo)
    
    # Guardar en archivo si se especifica ruta de salida
    if ruta_salida:
        os.makedirs(os.path.dirname(ruta_salida), exist_ok=True)
        with open(ruta_salida, "w", encoding="utf-8") as f:
            f.write(resultado)
        print(f"✓ Guardado: {ruta_salida}")
    
    return resultado


def procesar_todas_las_presentaciones(carpeta_base: str):
    """
    Busca y procesa todos los archivos .pptx en una carpeta y subcarpetas.
    """
    carpeta_base = Path(carpeta_base)
    carpeta_salida = carpeta_base / "presentaciones_texto"
    
    archivos_pptx = list(carpeta_base.rglob("*.pptx"))
    
    print(f"\n📁 Encontrados {len(archivos_pptx)} archivos PowerPoint\n")
    
    for archivo in archivos_pptx:
        try:
            # Crear nombre de archivo de salida
            nombre_relativo = archivo.relative_to(carpeta_base)
            ruta_txt = carpeta_salida / nombre_relativo.with_suffix(".txt")
            
            print(f"📄 Procesando: {archivo.name}")
            extraer_texto_pptx(str(archivo), str(ruta_txt))
            
        except Exception as e:
            print(f"❌ Error procesando {archivo.name}: {e}")
    
    print(f"\n✅ Proceso completado. Archivos guardados en: {carpeta_salida}")


if __name__ == "__main__":
    # Carpeta base del proyecto
    CARPETA_PROYECTO = r"c:\proyectos\WebSite-DigitalX\imagenes"
    
    procesar_todas_las_presentaciones(CARPETA_PROYECTO)
